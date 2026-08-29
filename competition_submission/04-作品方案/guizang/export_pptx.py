from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


HERE = Path(__file__).resolve().parent
SCREEN_DIR = HERE / "qa" / "screens"
OUTPUT = HERE.parent / "星火智学_作品方案_Guizang_DRAFT.pptx"


def main() -> None:
    images = sorted(SCREEN_DIR.glob("slide-*.png"))
    if len(images) != 12:
        raise SystemExit(
            "Expected 12 QA screenshots. Run: node qa_screenshots.mjs"
        )

    deck = Presentation()
    deck.slide_width = Inches(13.333333)
    deck.slide_height = Inches(7.5)
    deck.core_properties.title = "星火智学：本科刑法个性化学习与案件实训"
    deck.core_properties.subject = "XH-202620 比赛作品方案 Guizang视觉草案"
    deck.core_properties.author = "星火智学团队"
    deck.core_properties.comments = (
        "由Guizang网页PPT的1600x900视觉验收截图生成；"
        "网页源码保留动效与可编辑内容，本PPTX用于提交预览。"
    )

    blank = deck.slide_layouts[6]
    for image in images:
        slide = deck.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=deck.slide_width,
            height=deck.slide_height,
        )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    deck.save(OUTPUT)
    print(
        json.dumps(
            {
                "output": str(OUTPUT),
                "slides": len(deck.slides),
                "bytes": OUTPUT.stat().st_size,
                "under_100_mb": OUTPUT.stat().st_size < 100 * 1024 * 1024,
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
